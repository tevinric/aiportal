import streamlit as st
import pptx
import docx
import openpyxl
import Functions


from langchain_community.vectorstores import FAISS                  # --> (U002) For creating the vectorstore of the embeddings (text -> embed to vector -> store in vectorstore)  
from langchain_core.prompts import PromptTemplate        
from langchain_community.callbacks.manager import get_openai_callback
from langchain.chains.question_answering import load_qa_chain
from langchain_core.prompts import PromptTemplate
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.runnables import RunnablePassthrough

from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.prompts import PromptTemplate


@st.cache_resource
def get_vectorstore(text_content):
    # Create embeddings
    embeddings = AzureOpenAIEmbeddings(
        azure_endpoint=Functions.endpoint,
        api_key=Functions.api_key,
        api_version="2024-02-01",
        deployment="vectorai3",
    )
    
    # Split text into chunks
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=3000,
        chunk_overlap=100,
        length_function=len,
    )
    texts = text_splitter.split_text(text_content)
    
    # Create and return the vectorstore
    return FAISS.from_texts(texts, embeddings)


def extract_text_from_ppt(file): 
    from pptx import Presentation 
      
    presentation = Presentation(file)  
    text_content = []  
    
    for slide in presentation.slides:  
        for shape in slide.shapes:  
            if hasattr(shape, "text"):  
                text_content.append(shape.text)  
    return "\n".join(text_content)  


def summarize_text(client, deployment, temperature, text, summary_length, summary_type):  
    
    if summary_length =="Short":
        number_words = 250
    elif summary_length == "Medium":
        number_words = 500
    elif summary_length == "Long":
        number_words = 2000
    else:
        number_words = 5000
    
    response = client.chat.completions.create(  
        model=deployment,  
        messages=[  
            {"role": "system",
             "content": f"""You are an AI assistant designed to summarize documents. Your goal is to distill the information into clear summaries that highlight key points, actionable insights, and essential data.
                            When summarizing, consider the following guidelines:
                            Focus on Key Messages: Identify and articulate the main themes and objectives of the presentation.
                            Highlight Action Items: Clearly outline any recommended actions or decisions that need to be made.
                            Use Clear Language: Avoid jargon and technical terms unless necessary; aim for clarity and brevity.
                            Structure the Summary: Organize the summary into sections, such as Introduction, Key Findings, Recommendations, and Conclusion.
                            
                            Provide a summary in less than {number_words} words. 
                            The summary type must be {summary_type}. 
                            Concise means to summarise the points as is. 
                            Creative means to make the summary more creative with contextually related analogies that help the user understand the summary better. 
                            Narrtive means to generate the summary from the perspective of a story being told.
                            Non-technical means to create the summary for a non technical audience, i,e you must simplify technical jargon and specialised concepts. 
                            Technical mean to create a summary that retains technical concepts, ideas and content. 
                            Page by Page means to provide a summary of the content for each page. 
                            Section by Section means to provide a sumamry of the content for each section.
                            When provided with a PowerPoint presentation, generate a summary that meets these criteria, ensuring it is tailored for an executive audience."""                          
            },
            {"role": "user",
             "content": f"Please summarize the following text:\n\n{text}"
            }  
        ],
        temperature=temperature
    )  
    return response.choices[0].message.content


def extract_text_from_docx(file):
    import docx
    doc = docx.Document(file)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return "\n".join(text)

def extract_text_from_pdf(file):
    import fitz
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text.append(page.get_text())
    return "\n".join(text)

def extract_text_from_excel(file):
    import pandas as pd
    df = pd.read_excel(file)
    text = df.to_string(index=False)
    return text

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import json


def optimize_for_presentation(client, summary):
    """Optimize the summary content for presentation format using LLM"""
    response = client.chat.completions.create(
        model="gpt4omini",
        messages=[
            {"role": "system", 
             "content": """You are an expert at converting document summaries into presentation-ready content.
                          You must return valid JSON following this exact structure:
                          {
                              "title_slide": {
                                  "title": "Document Summary",
                                  "subtitle": "Key Points and Insights"
                              },
                              "agenda": ["Introduction", "Key Findings", "Recommendations"],
                              "sections": [
                                  {
                                      "title": "Section Title",
                                      "points": ["Point 1", "Point 2", "Point 3"]
                                  }
                              ],
                              "key_takeaways": ["Takeaway 1", "Takeaway 2", "Takeaway 3"]
                          }
                          
                          Guidelines:
                          - Break content into 3-7 points per section
                          - Use active voice and clear language
                          - Group related points together
                          - Do not include bullet points or dashes in the text
                          - Ensure all text is concise and presentation-friendly"""},
            {"role": "user",
             "content": f"Convert this summary into presentation-ready content, returning only the JSON structure:\n\n{summary}"}
        ],
        temperature=0.7
    )
    
    try:
        content = json.loads(response.choices[0].message.content)
    except json.JSONDecodeError:
        try:
            json_str = response.choices[0].message.content
            start = json_str.find('{')
            end = json_str.rfind('}') + 1
            if start != -1 and end != 0:
                content = json.loads(json_str[start:end])
            else:
                content = {
                    "title_slide": {
                        "title": "Document Summary",
                        "subtitle": "Key Points and Insights"
                    },
                    "agenda": ["Key Points", "Details", "Conclusion"],
                    "sections": [
                        {
                            "title": "Key Points",
                            "points": [point.lstrip('•- ').strip() for point in summary.split('\n')[:5]]
                        }
                    ],
                    "key_takeaways": [point.lstrip('•- ').strip() for point in summary.split('\n')[-3:]]
                }
        except Exception as e:
            content = {
                "title_slide": {
                    "title": "Document Summary",
                    "subtitle": "Key Points and Insights"
                },
                "agenda": ["Summary"],
                "sections": [
                    {
                        "title": "Summary",
                        "points": [summary]
                    }
                ],
                "key_takeaways": ["Please refer to the main summary for details"]
            }
    
    return content

def create_presentation(ppt_content):
    """Create a PowerPoint presentation from optimized content"""
    prs = Presentation()
    
    def add_text_to_shape(shape, text_list):
        """Helper function to add text to a shape with proper bullet formatting"""
        text_frame = shape.text_frame
        text_frame.clear()
        
        for idx, text in enumerate(text_list):
            paragraph = text_frame.add_paragraph()
            paragraph.text = text.strip()
            paragraph.font.size = Pt(24)
            # Only add bullet format if it's not the first paragraph in a title slide
            if not (shape.is_placeholder and shape.placeholder_format.idx == 0):
                paragraph.level = 0
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1] if len(title_slide.placeholders) > 1 else None
    
    title.text = ppt_content['title_slide']['title']
    if subtitle:
        subtitle.text = ppt_content['title_slide']['subtitle']
    
    # Agenda slide
    agenda_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = agenda_slide.shapes.title
    title.text = "Agenda"
    
    if len(agenda_slide.placeholders) > 1:
        body = agenda_slide.placeholders[1]
        add_text_to_shape(body, ppt_content['agenda'])
    else:
        body = agenda_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        add_text_to_shape(body, ppt_content['agenda'])
    
    # Content slides
    for section in ppt_content['sections']:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = section['title']
        
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            add_text_to_shape(body, section['points'])
        else:
            body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            add_text_to_shape(body, section['points'])
    
    # Key takeaways slide
    takeaway_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = takeaway_slide.shapes.title
    title.text = "Key Takeaways"
    
    if len(takeaway_slide.placeholders) > 1:
        body = takeaway_slide.placeholders[1]
        add_text_to_shape(body, ppt_content['key_takeaways'])
    else:
        body = takeaway_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        add_text_to_shape(body, ppt_content['key_takeaways'])
    
    # Save to BytesIO object
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    
    return pptx_buffer

def download_as_pptx(client, summary):
    """Main function to handle the conversion and download process"""
    # Optimize content for presentation
    optimized_content = optimize_for_presentation(client, summary)
    
    # Create presentation
    pptx_buffer = create_presentation(optimized_content)
    
    return pptx_buffer

def doc_summary(client):
    
    import tiktoken
    
    col1, col2, col3 = st.columns(3)
    
    with col1: 
        st.write(" ")
    with col2:
        st.image("Telesure-logo.png", width=200)
    with col3:
        st.write(" ")
        
    st.write(" ")
    
    summary_length = st.sidebar.selectbox("Summary Length", ["Short", "Medium", "Long", "Very Long"], 
                                        help="Short for < 100 words, Medium for < 500 words, Long for < 2000 words and Very Long for < 5000 words")
    summary_type = st.sidebar.selectbox("Summary Type", ["Concise", "Creative", "Executive", "Narrative", "Non-technical", "Technical", "Page by Page", "Section by Section"],
                                    help="Concise - Clear short summary points. Creative - Adds creative elements to summary. Narrtive - Summariese as a story narrtive.")
    
    summarizer_temperature = st.sidebar.slider("Temperature", 0.0, 1.0, 0.7, 
                                            help="Lower temperature to make the AI model more concise and less creative. Increase the tempature to make the AI model more creative")

    st.markdown("<h1 style='text-align: center;'>AI Document Summarization Tool</h1>", unsafe_allow_html=True)
    with st.expander("How to use"):
        st.write('''
            This tool provides two main functions:
            1. Document Summarization: Upload a document to get a summarized version.
            2. Document Chat: After uploading, you can ask questions about the document content.
            
            Supported document types: docx, pdf, xlsx, pptx
            
            Note: For large documents (20+ pages), process in smaller sections for better results.
        ''')
        
    uploaded_file = st.file_uploader("Choose a file")  
    
    if uploaded_file is not None:  
        file_type = uploaded_file.name.split('.')[-1]
        
        with st.spinner("Extracting text from the document..."):
            if file_type == "pptx":
                text_content = extract_text_from_ppt(uploaded_file)
            elif file_type == "docx":
                text_content = extract_text_from_docx(uploaded_file)
            elif file_type == "pdf":
                text_content = extract_text_from_pdf(uploaded_file)
            elif file_type == "xlsx":
                text_content = extract_text_from_excel(uploaded_file)
            else:
                text_content = None
        
        if text_content:
            # Create tabs for Summary and Chat
            summary_tab, chat_tab = st.tabs(["📝 Summary", "💬 Chat"])
            
            with summary_tab:
                if st.button("Summarize"):  
                    with st.spinner("Generating summary..."):  
                        summary = summarize_text(client, "gpt4omini", summarizer_temperature, 
                                            text_content, summary_length, summary_type)  
                        
                        # Store the summary in session state
                        st.session_state['current_summary'] = summary
                        
                    st.subheader("Summary:")   
                    st.write(summary)
                    
                    # Add download button
                    with st.spinner("Preparing PowerPoint download..."):
                        pptx_buffer = download_as_pptx(client, summary)
                        
                        st.download_button(
                            label="📥 Download as PowerPoint",
                            data=pptx_buffer.getvalue(),
                            file_name="summary_presentation.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key="download_pptx"
                        )
            
            with chat_tab:
                # Initialize chat history
                if "chat_history" not in st.session_state:
                    st.session_state.chat_history = []
                
                # Get or create vectorstore
                vectorstore = get_vectorstore(text_content)
                
                # Create chat interface
                st.write("Ask questions about your document:")
                
                # Display chat history
                for message in st.session_state.chat_history:
                    with st.chat_message(message["role"]):
                        st.markdown(message["content"])
                
                # Get user input
                if question := st.chat_input("Ask a question about your document"):
                    # Add user message to chat history
                    st.session_state.chat_history.append({"role": "user", "content": question})
                    
                    with st.chat_message("user", avatar="user.png"):
                        st.markdown(question)
                    
                    # Get relevant documents
                    docs = vectorstore.similarity_search(question)
                    
                    # Create QA chain
                    llm = AzureChatOpenAI(
                        azure_endpoint=Functions.endpoint,
                        api_key=Functions.api_key,
                        api_version="2024-02-01",
                        deployment_name="gpt4omini",
                        temperature=0.7
                    )
                    
                    prompt_template = """Use the following pieces of context to answer the question at the end. 
                    If you don't know the answer, just say that you don't know, don't try to make up an answer.

                    {context}

                    Question: {question}
                    Answer:"""
                    
                    prompt = PromptTemplate(
                        template=prompt_template,
                        input_variables=["context", "question"]
                    )

                    # Create the document chain
                    document_chain = create_stuff_documents_chain(
                        llm=llm,
                        prompt=prompt
                    )

                    # Create the retrieval chain
                    retrieval_chain = {
                        "context": lambda x: vectorstore.similarity_search(x["question"]),
                        "question": RunnablePassthrough()
                    }

                    # Combine the chains
                    chain = retrieval_chain | document_chain
                    
                    with st.chat_message("assistant", avatar="DNA Navigators.png"):
                        with st.spinner("Thinking..."):
                            # Get response from QA chain
                            answer = chain.invoke({"question": question})
                            st.markdown(answer)
                            
                            # Add assistant response to chat history
                            st.session_state.chat_history.append({"role": "assistant", "content": answer})                
        else:  
            st.warning("No text found in provided document")  

