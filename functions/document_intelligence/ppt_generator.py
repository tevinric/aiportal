import streamlit as st
from openai import AzureOpenAI
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import requests
from io import BytesIO
import os
from PIL import Image
import json
import base64

from config import api_key, endpoint


def initialize_azure_client():
    """Initialize Azure OpenAI client"""
    client = AzureOpenAI(
        api_key=api_key,  
        api_version="2024-02-15-preview",
        azure_endpoint=endpoint
    )
    return client

def generate_optimized_dalle_prompt(client, slide_content):
    """
    Generate an optimized prompt for DALLE-3 image generation using GPT-4
    to create detailed, professional, and contextually relevant prompts.
    """
    system_prompt = """You are an expert at creating optimal prompts for DALLE-3 image generation.
    Your role is to convert presentation slide content into detailed, specific image prompts that will
    generate professional, photorealistic images suitable for business presentations.

    Guidelines for prompt creation:
    1. Focus on photorealistic, professional imagery
    2. Avoid text or words in the image
    3. Emphasize lighting, composition, and atmosphere
    4. Use specific art direction terms
    5. Include clear style guidance
    6. Specify high-quality attributes
    7. Avoid abstract concepts unless necessary
    8. Ensure business-appropriate content

    Return only the prompt text without any explanations or additional content."""

    user_prompt = f"""Create a DALLE-3 prompt for a presentation slide with:
    Title: {slide_content['title']}
    Content: {slide_content['content']}
    Key Points: {', '.join(slide_content['key_points'])}
    
    The image should be suitable for a professional presentation."""

    try:
        response = client.chat.completions.create(
            model=st.session_state.model_deployment,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7,
            max_tokens=300
        )
        
        optimized_prompt = response.choices[0].message.content.strip()
        
        # Add consistent style guidelines to ensure professional quality
        style_suffix = (
            "Create this as a photorealistic image with professional lighting, "
            "shallow depth of field, high detail, 4K quality. "
            "Style: modern corporate photography, editorial quality. "
            "No text or words should appear in the image. "
            "Use natural lighting and professional composition."
        )
        
        return f"{optimized_prompt} {style_suffix}"
    except Exception as e:
        st.warning(f"Error generating optimized prompt: {str(e)}")
        return None

def generate_enhanced_dalle_image(client, slide_content, size="1024x1024"):
    """Generate enhanced images using optimized DALLE-3 prompts"""
    try:
        # Generate optimized prompt
        optimized_prompt = generate_optimized_dalle_prompt(client, slide_content)
        if not optimized_prompt:
            return None
            
        # Log the optimized prompt for debugging (optional)
        st.write(f"Generating image with optimized prompt: {optimized_prompt}")
        
        # Generate image using DALLE-3
        response = client.images.generate(
            model="Dalle3",  # Use your actual DALLE-3 deployment name
            prompt=optimized_prompt,
            n=1,
            size=size
        )
        
        # Get the image URL
        image_url = response.data[0].url
        
        # Download the image
        try:
            image_response = requests.get(image_url)
            if image_response.status_code == 200:
                return BytesIO(image_response.content)
            else:
                st.warning(f"Failed to download image: HTTP {image_response.status_code}")
                return None
        except requests.exceptions.RequestException as e:
            st.warning(f"Error downloading image: {str(e)}")
            return None
            
    except Exception as e:
        st.warning(f"Error generating image: {str(e)}")
        return None

def generate_presentation_images(client, slides):
    """Generate images for all presentation slides"""
    generated_images = []
    total_slides = len(slides)
    progress_bar = st.progress(0)
    
    for idx, slide in enumerate(slides):
        if slide.get("is_special", False):
            generated_images.append(None)
            continue
            
        with st.spinner(f"Generating image {idx + 1}/{total_slides}..."):
            try:
                image = generate_enhanced_dalle_image(client, slide)
                generated_images.append(image)
            except Exception as e:
                st.warning(f"Failed to generate image for slide {idx + 1}: {str(e)}")
                generated_images.append(None)
        
        progress = min(1.0, (idx + 1) / total_slides)
        progress_bar.progress(progress)
    
    return generated_images

def extract_text_from_pdf(pdf_file):
    """Extract text from PDF file"""
    text_content = []
    pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
    
    for page_num in range(pdf_document.page_count):
        page = pdf_document[page_num]
        text_content.append(page.get_text())
    
    return " ".join(text_content)

def generate_special_slides(client, content_data):
    """Generate comprehensive content for special slides using LLM"""
    system_prompt = """You are a professional presentation creator. Generate detailed content for special slides that ties the presentation together effectively."""

    special_slides = {}
    
    # Generate contents slide
    contents_prompt = f"""Create a detailed table of contents slide for the following presentation:
    Title: {content_data['title']}
    Subtitle: {content_data['subtitle']}
    Slides: {[{'title': slide['title'], 'key_points': slide['key_points']} for slide in content_data['slides']]}
    
    Return the response as a JSON object with:
    1. A brief introduction paragraph explaining what the presentation covers
    2. A list of key sections with brief descriptions
    3. Any special notes or highlights

    Format: {{
        "content": "introduction paragraph",
        "key_points": ["detailed point 1", "detailed point 2", ...]
    }}"""

    try:
        contents_response = client.chat.completions.create(
            model=st.session_state.model_deployment,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": contents_prompt}
            ],
            temperature=0.7,
            response_format={ "type": "json_object" }
        )
        
        contents_data = json.loads(contents_response.choices[0].message.content)
        special_slides['contents'] = {
            "title": "Contents",
            "content": contents_data['content'],
            "key_points": contents_data['key_points'],
            "is_special": True,
            "special_type": "contents"
        }
    except Exception as e:
        st.warning(f"Error generating contents: {str(e)}")
        special_slides['contents'] = {
            "title": "Contents",
            "content": "Presentation Overview",
            "key_points": [slide['title'] for slide in content_data['slides']],
            "is_special": True,
            "special_type": "contents"
        }

    # Generate conclusion slide
    conclusion_prompt = f"""Create a comprehensive conclusion for this presentation:
    Title: {content_data['title']}
    Subtitle: {content_data['subtitle']}
    Content: {json.dumps([{
        'title': slide['title'],
        'content': slide['content'],
        'key_points': slide['key_points']
    } for slide in content_data['slides']])}
    
    Generate a comprehensive conclusion that includes:
    1. A summary of the main presentation points
    2. Key takeaways and insights
    3. Final thoughts or recommendations
    4. Next steps or call to action
    
    Format the response as a JSON object:
    {{
        "content": "detailed summary paragraph",
        "key_points": [
            "key takeaway 1",
            "key takeaway 2",
            "recommendation",
            "call to action"
        ]
    }}"""

    try:
        conclusion_response = client.chat.completions.create(
            model=st.session_state.model_deployment,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": conclusion_prompt}
            ],
            temperature=0.7,
            response_format={ "type": "json_object" }
        )
        
        conclusion_data = json.loads(conclusion_response.choices[0].message.content)
        special_slides['conclusion'] = {
            "title": "Conclusion",
            "content": conclusion_data['content'],
            "key_points": conclusion_data['key_points'],
            "is_special": True,
            "special_type": "conclusion"
        }
    except Exception as e:
        st.warning(f"Error generating conclusion: {str(e)}")
        special_slides['conclusion'] = {
            "title": "Conclusion",
            "content": "Key takeaways from the presentation",
            "key_points": ["Summary point 1", "Summary point 2", "Summary point 3"],
            "is_special": True,
            "special_type": "conclusion"
        }

    return special_slides

def get_content_sections(client, text, style, num_slides, include_contents, include_conclusion, include_references, max_tokens=10000):
    """Use Azure OpenAI to analyze and structure the content"""
    system_prompt = f"""You are a professional presentation creator. Create a presentation outline with exactly {num_slides} content slides.
    You must respond with valid JSON only, using the following structure:
    {{
        "title": "Main presentation title",
        "subtitle": "Subtitle or tagline",
        "slides": [
            {{
                "title": "Slide title",
                "content": "Slide content (4-5 sentences)",
                "key_points": ["key point 1", "key point 2", "key point 3"]
            }}
        ]
    }}
    Break down the content into clear, concise points.
    Each slide should have a clear title, detailed content, and 3-4 key points."""

    try:
        # First generate main content
        response = client.chat.completions.create(
            model=st.session_state.model_deployment,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Create a {style} presentation from:\n{text[:max_tokens]}"}
            ],
            temperature=0.7,
            response_format={ "type": "json_object" }
        )
        
        content = json.loads(response.choices[0].message.content)
        
        # Generate special slides using the main content
        special_slides = generate_special_slides(client, content)
        
        # Construct final slides list
        final_slides = []
        
        # Add contents slide if selected
        if include_contents:
            final_slides.append(special_slides['contents'])
        
        # Add main content slides
        final_slides.extend(content["slides"])
        
        # Add conclusion slide if selected
        if include_conclusion:
            final_slides.append(special_slides['conclusion'])
        
        # Add references slide if selected
        if include_references:
            references_prompt = f"""Generate proper references and citations for this presentation:
            Title: {content['title']}
            Content: {text[:1000]}  # Using first 1000 chars for context
            
            Create 3-5 properly formatted references."""
            
            try:
                references_response = client.chat.completions.create(
                    model=st.session_state.model_deployment,
                    messages=[
                        {"role": "system", "content": "Generate academic references and citations."},
                        {"role": "user", "content": references_prompt}
                    ],
                    temperature=0.7
                )
                
                references = references_response.choices[0].message.content.split('\n')
                references = [ref.strip() for ref in references if ref.strip()]
                
                references_slide = {
                    "title": "References",
                    "content": "Sources and Citations",
                    "key_points": references,
                    "is_special": True,
                    "special_type": "references"
                }
            except Exception as e:
                st.warning(f"Error generating references: {str(e)}")
                references_slide = {
                    "title": "References",
                    "content": "Sources and Citations",
                    "key_points": ["Reference 1", "Reference 2", "Reference 3"],
                    "is_special": True,
                    "special_type": "references"
                }
            
            final_slides.append(references_slide)
        
        content["slides"] = final_slides
        return content
            
    except Exception as e:
        st.error(f"Error generating presentation structure: {str(e)}")
        return {
            "title": "Document Analysis",
            "subtitle": "Generated Summary",
            "slides": [
                {
                    "title": "Key Points",
                    "content": text[:200] + "...",
                    "key_points": ["Point 1", "Point 2", "Point 3"],
                    "is_special": False
                }
            ]
        }

def get_content_from_prompt(client, prompt, style, num_slides, include_contents, include_conclusion, include_references, max_tokens=10000):
    """Use Azure OpenAI to generate presentation content from a prompt"""
    system_prompt = f"""You are a professional presentation creator. Create a detailed {style} presentation with exactly {num_slides} content slides based on the given prompt.
    
    Requirements:
    1. Each slide must have 3-4 detailed key points
    2. Content should be informative and well-structured
    3. Language should match the {style} style
    4. Each slide's content should be 4-5 sentences
    
    You must respond with valid JSON only, using this structure:
    {{
        "title": "Main presentation title - make it engaging and relevant",
        "subtitle": "Informative subtitle that supports the title",
        "slides": [
            {{
                "title": "Clear and specific slide title",
                "content": "Detailed content explaining the slide topic (4-5 sentences)",
                "key_points": ["detailed key point 1", "detailed key point 2", "detailed key point 3"]
            }}
        ]
    }}"""

    try:
        # Generate main content from prompt
        response = client.chat.completions.create(
            model=st.session_state.model_deployment,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Create a detailed presentation about: {prompt}"}
            ],
            temperature=0.7,
            response_format={ "type": "json_object" }
        )
        
        content = json.loads(response.choices[0].message.content)
        
        # Generate special slides using the main content
        special_slides = generate_special_slides(client, content)
        
        # Construct final slides list
        final_slides = []
        
        # Add contents slide if selected
        if include_contents:
            contents_slide = special_slides['contents']
            contents_slide["is_special"] = True
            contents_slide["special_type"] = "contents"
            final_slides.append(contents_slide)
        
        # Add main content slides
        for slide in content["slides"]:
            slide["is_special"] = False
            final_slides.append(slide)
        
        # Add conclusion slide if selected
        if include_conclusion:
            conclusion_slide = special_slides['conclusion']
            conclusion_slide["is_special"] = True
            conclusion_slide["special_type"] = "conclusion"
            final_slides.append(conclusion_slide)
        
        # Add references slide if selected
        if include_references:
            references_prompt = f"""Generate academic references and citations for this presentation topic:
            Title: {content['title']}
            Topic: {prompt}
            
            Create 3-5 properly formatted references."""
            
            try:
                references_response = client.chat.completions.create(
                    model=st.session_state.model_deployment,
                    messages=[
                        {"role": "system", "content": "Generate academic references and citations."},
                        {"role": "user", "content": references_prompt}
                    ],
                    temperature=0.7
                )
                
                references = references_response.choices[0].message.content.split('\n')
                references = [ref.strip() for ref in references if ref.strip()]
                
                references_slide = {
                    "title": "References",
                    "content": "Sources and Citations",
                    "key_points": references,
                    "is_special": True,
                    "special_type": "references"
                }
                final_slides.append(references_slide)
                
            except Exception as e:
                st.warning(f"Error generating references: {str(e)}")
        
        content["slides"] = final_slides
        return content
            
    except Exception as e:
        st.error(f"Error generating presentation structure: {str(e)}")
        return None

def apply_theme_colors(theme):
    """Return extended color scheme based on selected theme"""
    color_schemes = {
        "Professional": {
            "background": RGBColor(255, 255, 255),
            "title": RGBColor(31, 73, 125),
            "text": RGBColor(0, 0, 0),
            "accent": RGBColor(0, 112, 192),
            "footer": RGBColor(89, 89, 89),
            "bullet": RGBColor(31, 73, 125)
        },
        "Modern": {
            "background": RGBColor(240, 240, 240),
            "title": RGBColor(41, 128, 185),
            "text": RGBColor(44, 62, 80),
            "accent": RGBColor(52, 152, 219),
            "footer": RGBColor(149, 165, 166),
            "bullet": RGBColor(41, 128, 185)
        },
        "Creative": {
            "background": RGBColor(255, 250, 240),
            "title": RGBColor(230, 126, 34),
            "text": RGBColor(44, 62, 80),
            "accent": RGBColor(243, 156, 18),
            "footer": RGBColor(149, 165, 166),
            "bullet": RGBColor(230, 126, 34)
        }
    }
    return color_schemes.get(theme, color_schemes["Professional"])

def add_slide_number(slide, slide_number, total_slides, colors):
    """Add slide number to the slide"""
    footer_shape = slide.shapes.add_textbox(
        left=Inches(12),
        top=Inches(7),
        width=Inches(1),
        height=Inches(0.3)
    )
    footer_text_frame = footer_shape.text_frame
    p = footer_text_frame.paragraphs[0]
    p.text = f"{slide_number}/{total_slides}"
    p.font.size = Pt(12)
    p.font.name = "Arial"
    p.font.color.rgb = colors["footer"]
    p.alignment = PP_ALIGN.RIGHT

def apply_content_formatting(paragraph, colors, is_bullet=False):
    """Apply consistent formatting to content paragraphs"""
    paragraph.font.size = Pt(14)
    paragraph.font.name = "Arial"
    paragraph.font.color.rgb = colors["text"]
    
    if is_bullet:
        paragraph.level = 0
        paragraph.bullet = True
        # Format bullet points
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.font.color.rgb = colors["bullet"]
    
    # Add proper line spacing
    paragraph.line_spacing = 1.15
    # Add space after paragraph
    paragraph.space_after = Pt(12)

def add_image_to_slide(slide, image_stream, layout_style):
    """Add image to slide with proper positioning"""
    if not image_stream:
        return
        
    try:
        if layout_style == "Content with Image":
            # Position image on the right side
            left = Inches(7)
            top = Inches(1.5)
            width = Inches(5.5)
            height = Inches(5)
        else:
            # Position image below content
            left = Inches(1)
            top = Inches(4)
            width = Inches(5)
            height = Inches(3)
        
        # Add image to slide
        slide.shapes.add_picture(
            image_stream,
            left,
            top,
            width,
            height
        )
    except Exception as e:
        st.warning(f"Error adding image to slide: {str(e)}")

def create_presentation(content_data, theme, layout_style, include_images=True):
    """Create PowerPoint presentation using python-pptx with enhanced formatting and images"""
    prs = Presentation()
    colors = apply_theme_colors(theme)
    
    # Set slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    total_slides = len(content_data["slides"]) + 1
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Format title
    title = slide.shapes.title
    title.top = Inches(2)
    title.left = Inches(1)
    title.width = Inches(11.333)
    title.height = Inches(1.5)
    
    # Format subtitle
    subtitle = slide.placeholders[1]
    subtitle.top = Inches(3.75)
    subtitle.left = Inches(1)
    subtitle.width = Inches(11.333)
    
    # Set title and subtitle text with formatting
    title.text = content_data["title"]
    title_para = title.text_frame.paragraphs[0]
    title_para.font.name = "Arial"
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = colors["title"]
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle.text = content_data["subtitle"]
    subtitle_para = subtitle.text_frame.paragraphs[0]
    subtitle_para.font.name = "Arial"
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = colors["text"]
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add slide number
    add_slide_number(slide, 1, total_slides, colors)
    
    # Content slides
    for idx, slide_data in enumerate(content_data["slides"], start=1):
        # Choose layout based on slide type and whether images are included
        if slide_data.get("is_special", False):
            layout = prs.slide_layouts[1]  # Title and Content layout for special slides
        else:
            layout = prs.slide_layouts[1]  # Using Title and Content layout for all slides
        
        slide = prs.slides.add_slide(layout)
        
        # Add title and content formatting
        title_shape = slide.shapes.title
        title_shape.top = Inches(0.4)
        title_shape.left = Inches(0.5)
        title_shape.width = Inches(12.33)
        title_shape.height = Inches(0.8)
        
        title_shape.text = slide_data["title"]
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.name = "Arial"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = colors["title"]
        title_para.alignment = PP_ALIGN.LEFT
        
        # Add accent line
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            Inches(1.3),
            Inches(12.33),
            Inches(0.03)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = colors["accent"]
        accent_line.line.fill.background()
        
        # Determine content width based on layout style and image presence
        if layout_style in ["Content with Image", "Bullets with Image"] and not slide_data.get("is_special", False):
            content_width = Inches(6.5)
        else:
            content_width = Inches(11.93)
            
        content_box = slide.shapes.add_textbox(
            left=Inches(0.7),
            top=Inches(1.4),
            width=content_width,
            height=Inches(5.2)
        )
        
        # Add content with formatting
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        # Handle different layout styles
        if layout_style in ["Bullet Points", "Bullets with Image"]:
            # Add main content as a paragraph if it exists and isn't empty
            if slide_data.get("content") and str(slide_data["content"]).strip():
                p = text_frame.add_paragraph()
                p.text = str(slide_data["content"])
                apply_content_formatting(p, colors, False)
                p.space_after = Pt(20)  # Add space between content and bullets
            
            # Add bullet points
            if "key_points" in slide_data and slide_data["key_points"]:
                for point in slide_data["key_points"]:
                    p = text_frame.add_paragraph()
                    p.text = str(point)  # Convert to string to handle any non-string points
                    apply_content_formatting(p, colors, True)
        else:
            # Standard content layout
            p = text_frame.add_paragraph()
            p.text = str(slide_data.get("content", ""))
            apply_content_formatting(p, colors, False)
        
        # Add image if available and layout allows
        if include_images and "image" in slide_data and slide_data["image"] and not slide_data.get("is_special", False):
            if layout_style == "Content with Image":
                # Position image on the right side
                left = Inches(7)
                top = Inches(1.5)
                width = Inches(5.5)
                height = Inches(5)
            elif layout_style == "Bullets with Image":
                # Position image below the bullet points
                left = Inches(0.7)
                top = Inches(4)
                width = Inches(6.5)
                height = Inches(3)
            else:
                # Default image position below content
                left = Inches(1)
                top = Inches(4)
                width = Inches(5)
                height = Inches(3)
            
            try:
                slide.shapes.add_picture(
                    slide_data["image"],
                    left,
                    top,
                    width,
                    height
                )
            except Exception as e:
                st.warning(f"Error adding image to slide {idx}: {str(e)}")
        
        add_slide_number(slide, idx + 1, total_slides, colors)
    
    return prs

def regenerate_single_image(client, slide_data):
    """Regenerate image for a single slide"""
    try:
        return generate_enhanced_dalle_image(client, slide_data)
    except Exception as e:
        st.warning(f"Error regenerating image: {str(e)}")
        return None



def ppt_app():
    
    # Initialize session state variables
    if "client" not in st.session_state:
        st.session_state.client = initialize_azure_client()
    if "content_data" not in st.session_state:
        st.session_state.content_data = None
    if "current_step" not in st.session_state:
        st.session_state.current_step = "select_mode"
    if "mode" not in st.session_state:
        st.session_state.mode = None
    if "image_descriptions" not in st.session_state:
        st.session_state.image_descriptions = []
    if "generated_images" not in st.session_state:
        st.session_state.generated_images = []
    if "regenerated_images" not in st.session_state:
        st.session_state.regenerated_images = {}

    st.title("AI PowerPoint Presentation Creator")
    
    # Sidebar Configuration
    st.sidebar.title("Presentation Settings")
    
    # Model Selection
    model_deployments = ["gpt4o", "gpt4omini"]
    st.session_state.model_deployment = st.sidebar.selectbox(
        "Select GPT Model",
        model_deployments,
        index=0
    )
    
    # Configuration Options
    st.sidebar.subheader("Configuration")
    num_slides = st.sidebar.number_input("Number of Content Slides", min_value=3, value=5)

    # Visual Options
    st.sidebar.subheader("Visual Settings")
    include_images = st.sidebar.toggle("Generate Images", value=True)
    
    # Special Slides Options
    st.sidebar.subheader("Special Slides")
    include_contents = st.sidebar.toggle("Include Contents Slide", value=True)
    include_conclusion = st.sidebar.toggle("Include Conclusion Slide", value=True)
    include_references = st.sidebar.toggle("Include References Slide", value=True)
    
    # Theme and Layout Settings
    st.sidebar.subheader("Design Settings")
    theme = st.sidebar.selectbox(
        "Select Theme",
        ["Professional", "Modern", "Creative"],
        index=0
    )
    
    layout_style = st.sidebar.selectbox(
        "Select Layout Style",
        ["Content with Image", "Bullet Points", "Bullets with Image", "Minimal"],
        index=0
    )
    
    presentation_style = st.sidebar.selectbox(
        "Presentation Style",
        ["Business", "Academic", "Creative", "Technical"],
        index=0
    )

    # Step 1: Select Mode
    if st.session_state.current_step == "select_mode":
        st.write("Choose how you want to create your presentation:")
        col1, col2 = st.columns(2)
        
        with col1:
            if st.button("Create from Text Prompt", use_container_width=True):
                st.session_state.mode = "prompt"
                st.session_state.current_step = "input"
                st.rerun()
                
        with col2:
            if st.button("Create from PDF Document", use_container_width=True):
                st.session_state.mode = "document"
                st.session_state.current_step = "input"
                st.rerun()

    # Step 2: Input (either prompt or document)
    elif st.session_state.current_step == "input":
        if st.session_state.mode == "prompt":
            st.write("Enter your presentation topic or description:")
            prompt = st.text_area(
                "Presentation Prompt", 
                height=150,
                placeholder="Describe what you want your presentation to be about in detail. The more specific you are, the better the results will be."
            )
            
            if prompt and st.button("Generate Presentation"):
                with st.spinner("Generating presentation content..."):
                    st.session_state.content_data = get_content_from_prompt(
                        st.session_state.client,
                        prompt,
                        presentation_style,
                        num_slides,
                        include_contents,
                        include_conclusion,
                        include_references
                    )
                    
                    if st.session_state.content_data:
                        if include_images:
                            st.session_state.current_step = "generate_images"
                        else:
                            st.session_state.current_step = "create_presentation"
                        st.rerun()
                    else:
                        st.error("Failed to generate presentation content. Please try again with a more detailed prompt.")
                    
        else:  # document mode
            st.write("Upload your PDF document:")
            uploaded_file = st.file_uploader("Choose a PDF file", type="pdf")
            if uploaded_file is not None:
                with st.spinner("Processing PDF and generating presentation content..."):
                    text_content = extract_text_from_pdf(uploaded_file)
                    st.session_state.content_data = get_content_sections(
                        st.session_state.client,
                        text_content,
                        presentation_style,
                        num_slides,
                        include_contents,
                        include_conclusion,
                        include_references
                    )
                    
                    if include_images:
                        st.session_state.current_step = "generate_images"
                    else:
                        st.session_state.current_step = "create_presentation"
                    st.rerun()

    # Step 3: Generate Images
    elif st.session_state.current_step == "generate_images":
        st.subheader("Generate Images")
        
        if not st.session_state.content_data:
            st.error("No content data found. Please start over.")
            if st.button("Start Over"):
                st.session_state.current_step = "select_mode"
                st.rerun()
            return
        
        # Generate images using enhanced function
        st.session_state.generated_images = generate_presentation_images(
            st.session_state.client,
            st.session_state.content_data["slides"]
        )
        
        st.session_state.current_step = "create_presentation"
        st.rerun()

    # Step 4: Create Presentation
    elif st.session_state.current_step == "create_presentation":
        st.subheader("Create Presentation")
        
        if not st.session_state.content_data:
            st.error("No content data found. Please start over.")
            if st.button("Start Over"):
                st.session_state.current_step = "select_mode"
                st.rerun()
            return
            
        presentation_data = {
            "title": st.session_state.content_data["title"],
            "subtitle": st.session_state.content_data["subtitle"],
            "slides": []
        }
        
        for idx, slide in enumerate(st.session_state.content_data["slides"]):
            slide_data = slide.copy()
            if include_images:
                # Use regenerated image if available, otherwise use original
                if idx in st.session_state.regenerated_images:
                    slide_data["image"] = st.session_state.regenerated_images[idx]
                elif idx < len(st.session_state.generated_images):
                    slide_data["image"] = st.session_state.generated_images[idx]
            presentation_data["slides"].append(slide_data)

        with st.spinner("Creating presentation..."):
            # Create a column layout for the preview and download section
            preview_col, download_col = st.columns([2, 1])
            
            # Preview column
            with preview_col:
                st.subheader("Content Preview")
                for idx, slide in enumerate(presentation_data["slides"]):
                    with st.expander(f"Slide {idx + 1}: {slide['title']}", expanded=False):
                        # Create two columns for content and image controls
                        content_col, image_col = st.columns([2, 1])
                        
                        with content_col:
                            st.markdown("**Content:**")
                            st.write(slide['content'])
                            
                            st.markdown("**Key Points:**")
                            for point in slide['key_points']:
                                st.markdown(f"• {point}")
                        

                            if include_images and not slide.get("is_special", False):
                                if "image" in slide and slide["image"]:
                                    st.markdown("**Generated Image:**")
                                    # Display image at reduced size
                                    st.image(slide["image"], 
                                            caption=f"Generated image for slide {idx + 1}",
                                            width=300)  # Reduced width for preview
                                    
                                    # Add regenerate button for this specific slide
                                    if st.button(f"🔄 Regenerate Image", key=f"regen_{idx}"):
                                        with st.spinner("Regenerating image..."):
                                            new_image = regenerate_single_image(
                                                st.session_state.client,
                                                slide
                                            )
                                            if new_image:
                                                st.session_state.regenerated_images[idx] = new_image
                                                # Update the presentation data
                                                presentation_data["slides"][idx]["image"] = new_image
                                                st.rerun()
            
            # Generate the presentation with latest images
            presentation = create_presentation(
                presentation_data,
                theme,
                layout_style,
                include_images
            )
            
            pptx_io = BytesIO()
            presentation.save(pptx_io)
            pptx_io.seek(0)
            
            # Download column
            with download_col:
                st.markdown("### Download Options")
                st.download_button(
                    label="📥 Download PowerPoint",
                    data=pptx_io,
                    file_name="generated_presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
                
                # Add options for additional formats or variations
                st.markdown("### Additional Options")
                if st.button("Create New Presentation", use_container_width=True):
                    for key in ["content_data", "current_step", "mode", "image_descriptions", 
                              "generated_images", "regenerated_images"]:
                        if key in st.session_state:
                            del st.session_state[key]
                    st.rerun()
                
                if st.button("Regenerate All Images", use_container_width=True, disabled=not include_images):
                    st.session_state.current_step = "generate_images"
                    if "generated_images" in st.session_state:
                        del st.session_state.generated_images
                    if "regenerated_images" in st.session_state:
                        del st.session_state.regenerated_images
                    st.rerun()
                
                st.markdown("---")
                st.markdown("""
                ### Presentation Stats
                - Total Slides: {}
                - Images Generated: {}
                - Theme: {}
                - Layout: {}
                """.format(
                    len(presentation_data["slides"]),
                    sum(1 for slide in presentation_data["slides"] if "image" in slide and slide["image"]),
                    theme,
                    layout_style
                ))
    
    # Add a "Start Over" button in the sidebar that's always available
    if st.sidebar.button("Start Over", use_container_width=True):
        for key in ["content_data", "current_step", "mode", "image_descriptions", 
                   "generated_images", "regenerated_images"]:
            if key in st.session_state:
                del st.session_state[key]
        st.rerun()

    # Add footer with helpful information
    st.sidebar.markdown("---")
    st.sidebar.markdown("""
    ### Tips for Best Results:
    - Be specific in your presentation topic or PDF content
    - Choose a theme that matches your content style
    - Review all slides before downloading
    - Use image generation for key concepts
    - Consider regenerating images if needed
    
    ### Image Generation Tips:
    - Images work best with concrete concepts
    - Abstract ideas may need multiple attempts
    - Professional themes produce better results
    - Review generated images for relevance
    """)
