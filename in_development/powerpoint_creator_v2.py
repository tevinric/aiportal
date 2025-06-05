import streamlit as st
from openai import AzureOpenAI
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import requests
from io import BytesIO
import os
from PIL import Image
import json
from dotenv import load_dotenv
import base64

load_dotenv()

def initialize_azure_client():
    """Initialize Azure OpenAI client"""
    client = AzureOpenAI(
        api_key=os.getenv("api_key"),  
        api_version="2024-02-15-preview",
        azure_endpoint=os.getenv("endpoint")
    )
    return client

def extract_text_from_pdf(pdf_file):
    """Extract text from PDF file"""
    text_content = []
    pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
    
    for page_num in range(pdf_document.page_count):
        page = pdf_document[page_num]
        text_content.append(page.get_text())
    
    return " ".join(text_content)
def generate_special_slides(client, content_data):
    """Generate comprehensive content for special slides using LLM"""
    system_prompt = """You are a professional presentation creator. Generate detailed content for special slides that ties the presentation together effectively."""

    special_slides = {}
    
    # Generate contents slide
    contents_prompt = f"""Create a detailed table of contents slide for the following presentation:
    Title: {content_data['title']}
    Subtitle: {content_data['subtitle']}
    Slides: {[{'title': slide['title'], 'key_points': slide['key_points']} for slide in content_data['slides']]}
    
    Return the response as a JSON object with:
    1. A brief introduction paragraph explaining what the presentation covers
    2. A list of key sections with brief descriptions
    3. Any special notes or highlights

    Format: {{
        "content": "introduction paragraph",
        "key_points": ["detailed point 1", "detailed point 2", ...]
    }}"""

    try:
        contents_response = client.chat.completions.create(
            model=st.session_state.model_deployment,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": contents_prompt}
            ],
            temperature=0.7,
            response_format={ "type": "json_object" }
        )
        
        contents_data = json.loads(contents_response.choices[0].message.content)
        special_slides['contents'] = {
            "title": "Contents",
            "content": contents_data['content'],
            "key_points": contents_data['key_points'],
            "is_special": True,
            "special_type": "contents"
        }
    except Exception as e:
        st.warning(f"Error generating contents: {str(e)}")
        special_slides['contents'] = {
            "title": "Contents",
            "content": "Presentation Overview",
            "key_points": [slide['title'] for slide in content_data['slides']],
            "is_special": True,
            "special_type": "contents"
        }

    # Generate conclusion slide
    conclusion_prompt = f"""Create a comprehensive conclusion for this presentation:
    Title: {content_data['title']}
    Subtitle: {content_data['subtitle']}
    Content: {json.dumps([{
        'title': slide['title'],
        'content': slide['content'],
        'key_points': slide['key_points']
    } for slide in content_data['slides']])}
    
    Generate a comprehensive conclusion that includes:
    1. A summary of the main presentation points
    2. Key takeaways and insights
    3. Final thoughts or recommendations
    4. Next steps or call to action
    
    Format the response as a JSON object:
    {{
        "content": "detailed summary paragraph",
        "key_points": [
            "key takeaway 1",
            "key takeaway 2",
            "recommendation",
            "call to action"
        ]
    }}"""

    try:
        conclusion_response = client.chat.completions.create(
            model=st.session_state.model_deployment,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": conclusion_prompt}
            ],
            temperature=0.7,
            response_format={ "type": "json_object" }
        )
        
        conclusion_data = json.loads(conclusion_response.choices[0].message.content)
        special_slides['conclusion'] = {
            "title": "Conclusion",
            "content": conclusion_data['content'],
            "key_points": conclusion_data['key_points'],
            "is_special": True,
            "special_type": "conclusion"
        }
    except Exception as e:
        st.warning(f"Error generating conclusion: {str(e)}")
        special_slides['conclusion'] = {
            "title": "Conclusion",
            "content": "Key takeaways from the presentation",
            "key_points": ["Summary point 1", "Summary point 2", "Summary point 3"],
            "is_special": True,
            "special_type": "conclusion"
        }

    return special_slides

def get_content_sections(client, text, style, num_slides, include_contents, include_conclusion, include_references, max_tokens=10000):
    """Use Azure OpenAI to analyze and structure the content"""
    system_prompt = f"""You are a professional presentation creator. Create a presentation outline with exactly {num_slides} content slides.
    You must respond with valid JSON only, using the following structure:
    {{
        "title": "Main presentation title",
        "subtitle": "Subtitle or tagline",
        "slides": [
            {{
                "title": "Slide title",
                "content": "Slide content (4-5 sentences)",
                "key_points": ["key point 1", "key point 2", "key point 3"]
            }}
        ]
    }}
    Break down the content into clear, concise points.
    Each slide should have a clear title, detailed content, and 3-4 key points."""

    try:
        # First generate main content
        response = client.chat.completions.create(
            model=st.session_state.model_deployment,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Create a {style} presentation from:\n{text[:max_tokens]}"}
            ],
            temperature=0.7,
            response_format={ "type": "json_object" }
        )
        
        content = json.loads(response.choices[0].message.content)
        
        # Generate special slides using the main content
        special_slides = generate_special_slides(client, content)
        
        # Construct final slides list
        final_slides = []
        
        # Add contents slide if selected
        if include_contents:
            final_slides.append(special_slides['contents'])
        
        # Add main content slides
        final_slides.extend(content["slides"])
        
        # Add conclusion slide if selected
        if include_conclusion:
            final_slides.append(special_slides['conclusion'])
        
        # Add references slide if selected
        if include_references:
            references_prompt = f"""Generate proper references and citations for this presentation:
            Title: {content['title']}
            Content: {text[:1000]}  # Using first 1000 chars for context
            
            Create 3-5 properly formatted references."""
            
            try:
                references_response = client.chat.completions.create(
                    model=st.session_state.model_deployment,
                    messages=[
                        {"role": "system", "content": "Generate academic references and citations."},
                        {"role": "user", "content": references_prompt}
                    ],
                    temperature=0.7
                )
                
                references = references_response.choices[0].message.content.split('\n')
                references = [ref.strip() for ref in references if ref.strip()]
                
                references_slide = {
                    "title": "References",
                    "content": "Sources and Citations",
                    "key_points": references,
                    "is_special": True,
                    "special_type": "references"
                }
            except Exception as e:
                st.warning(f"Error generating references: {str(e)}")
                references_slide = {
                    "title": "References",
                    "content": "Sources and Citations",
                    "key_points": ["Reference 1", "Reference 2", "Reference 3"],
                    "is_special": True,
                    "special_type": "references"
                }
            
            final_slides.append(references_slide)
        
        content["slides"] = final_slides
        return content
            
    except Exception as e:
        st.error(f"Error generating presentation structure: {str(e)}")
        return {
            "title": "Document Analysis",
            "subtitle": "Generated Summary",
            "slides": [
                {
                    "title": "Key Points",
                    "content": text[:200] + "...",
                    "key_points": ["Point 1", "Point 2", "Point 3"],
                    "is_special": False
                }
            ]
        }
        
def generate_image_descriptions(client, slides):
    """Generate image descriptions based on slide content"""
    system_prompt = """You are a professional presentation designer. Generate appropriate image descriptions 
    for presentation slides based on their content. The descriptions should be specific and detailed enough 
    for an AI image generator to create relevant, professional visuals."""

    descriptions = []
    for slide in slides:
        user_prompt = f"""Create an image description for a presentation slide with the following:
        Title: {slide['title']}
        Content: {slide['content']}
        Key Points: {', '.join(slide['key_points'])}
        
        The image should be professional and suitable for a business presentation."""

        try:
            response = client.chat.completions.create(
                model=st.session_state.model_deployment,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.7
            )
            
            descriptions.append(response.choices[0].message.content)
        except Exception as e:
            descriptions.append("A professional business illustration representing the slide content")
            st.warning(f"Error generating image description: {str(e)}")

    return descriptions

def generate_dalle_image(client, prompt, size="1024x1024"):
    """Generate image using DALLE-3"""
    try:
        response = client.images.generate(
            model="Dalle3", # Use your actual DALLE-3 deployment name
            prompt=f"Create a professional presentation slide image: {prompt}. Style: clean, professional, suitable for business presentation.",
            n=1,
            size=size
        )
        
        # Get the image URL
        image_url = response.data[0].url
        
        # Download the image
        try:
            image_response = requests.get(image_url)
            if image_response.status_code == 200:
                return BytesIO(image_response.content)
            else:
                st.warning(f"Failed to download image: HTTP {image_response.status_code}")
                return None
        except requests.exceptions.RequestException as e:
            st.warning(f"Error downloading image: {str(e)}")
            return None
            
    except Exception as e:
        st.warning(f"Error generating image: {str(e)}")
        return None
    
def apply_theme_colors(theme):
    """Return extended color scheme based on selected theme"""
    color_schemes = {
        "Professional": {
            "background": RGBColor(255, 255, 255),
            "title": RGBColor(31, 73, 125),
            "text": RGBColor(0, 0, 0),
            "accent": RGBColor(0, 112, 192),
            "footer": RGBColor(89, 89, 89),
            "bullet": RGBColor(31, 73, 125)
        },
        "Modern": {
            "background": RGBColor(240, 240, 240),
            "title": RGBColor(41, 128, 185),
            "text": RGBColor(44, 62, 80),
            "accent": RGBColor(52, 152, 219),
            "footer": RGBColor(149, 165, 166),
            "bullet": RGBColor(41, 128, 185)
        },
        "Creative": {
            "background": RGBColor(255, 250, 240),
            "title": RGBColor(230, 126, 34),
            "text": RGBColor(44, 62, 80),
            "accent": RGBColor(243, 156, 18),
            "footer": RGBColor(149, 165, 166),
            "bullet": RGBColor(230, 126, 34)
        }
    }
    return color_schemes.get(theme, color_schemes["Professional"])

def add_slide_number(slide, slide_number, total_slides, colors):
    """Add slide number to the slide"""
    footer_shape = slide.shapes.add_textbox(
        left=Inches(12),
        top=Inches(7),
        width=Inches(1),
        height=Inches(0.3)
    )
    footer_text_frame = footer_shape.text_frame
    p = footer_text_frame.paragraphs[0]
    p.text = f"{slide_number}/{total_slides}"
    p.font.size = Pt(12)
    p.font.name = "Arial"
    p.font.color.rgb = colors["footer"]
    p.alignment = PP_ALIGN.RIGHT

def apply_content_formatting(paragraph, colors, is_bullet=False):
    """Apply consistent formatting to content paragraphs"""
    paragraph.font.size = Pt(14)
    paragraph.font.name = "Arial"
    paragraph.font.color.rgb = colors["text"]
    
    if is_bullet:
        paragraph.level = 0
        paragraph.bullet = True
        # Format bullet points
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.font.color.rgb = colors["bullet"]
    
    # Add proper line spacing
    paragraph.line_spacing = 1.15
    # Add space after paragraph
    paragraph.space_after = Pt(12)

def add_image_to_slide(slide, image_stream, layout_style):
    """Add image to slide with proper positioning"""
    if not image_stream:
        return
        
    try:
        if layout_style == "Content with Image":
            # Position image on the right side
            left = Inches(7)
            top = Inches(1.5)
            width = Inches(5.5)
            height = Inches(5)
        else:
            # Position image below content
            left = Inches(1)
            top = Inches(4)
            width = Inches(5)
            height = Inches(3)
        
        # Add image to slide
        slide.shapes.add_picture(
            image_stream,
            left,
            top,
            width,
            height
        )
    except Exception as e:
        st.warning(f"Error adding image to slide: {str(e)}")
        
def create_presentation(content_data, theme, layout_style, include_images=True):
    """Create PowerPoint presentation using python-pptx with enhanced formatting and images"""
    prs = Presentation()
    colors = apply_theme_colors(theme)
    
    # Set slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    total_slides = len(content_data["slides"]) + 1
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Format title
    title = slide.shapes.title
    title.top = Inches(2)
    title.left = Inches(1)
    title.width = Inches(11.333)
    title.height = Inches(1.5)
    
    # Format subtitle
    subtitle = slide.placeholders[1]
    subtitle.top = Inches(3.75)
    subtitle.left = Inches(1)
    subtitle.width = Inches(11.333)
    
    # Set title and subtitle text with formatting
    title.text = content_data["title"]
    title_para = title.text_frame.paragraphs[0]
    title_para.font.name = "Arial"
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = colors["title"]
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle.text = content_data["subtitle"]
    subtitle_para = subtitle.text_frame.paragraphs[0]
    subtitle_para.font.name = "Arial"
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = colors["text"]
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add slide number
    add_slide_number(slide, 1, total_slides, colors)
    
    # Content slides
    for idx, slide_data in enumerate(content_data["slides"], start=1):
        # Choose layout based on slide type and whether images are included
        if slide_data.get("is_special", False):
            layout = prs.slide_layouts[1]  # Title and Content layout for special slides
        else:
            layout = prs.slide_layouts[8] if (include_images and layout_style == "Content with Image") else prs.slide_layouts[1]
        
        slide = prs.slides.add_slide(layout)
        
        # Add title and content formatting
        title_shape = slide.shapes.title
        title_shape.top = Inches(0.4)
        title_shape.left = Inches(0.5)
        title_shape.width = Inches(12.33)
        title_shape.height = Inches(0.8)
        
        title_shape.text = slide_data["title"]
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.name = "Arial"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = colors["title"]
        title_para.alignment = PP_ALIGN.LEFT
        
        # Add accent line
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            Inches(1.3),
            Inches(12.33),
            Inches(0.03)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = colors["accent"]
        accent_line.line.fill.background()
        
        # Adjust content box width based on layout
        content_width = Inches(6.5) if layout_style == "Content with Image" else Inches(11.93)
        content_box = slide.shapes.add_textbox(
            left=Inches(0.7),
            top=Inches(1.4),
            width=content_width,
            height=Inches(5.2)
        )
        
        # Add content with formatting
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        # Add main content
        if layout_style == "Bullet Points":
            for point in slide_data["key_points"]:
                p = text_frame.add_paragraph()
                p.text = point
                apply_content_formatting(p, colors, True)
        else:
            p = text_frame.add_paragraph()
            p.text = slide_data["content"]
            apply_content_formatting(p, colors, False)
        
        # Add image if available
        if include_images and "image" in slide_data and slide_data["image"] and not slide_data.get("is_special", False):
            add_image_to_slide(slide, slide_data["image"], layout_style)
        
        add_slide_number(slide, idx + 1, total_slides, colors)
    
    return prs





def main():
    st.set_page_config(page_title="PDF to PowerPoint Creator", layout="wide")
    
    # Initialize session state variables
    if "client" not in st.session_state:
        st.session_state.client = initialize_azure_client()
    if "content_data" not in st.session_state:
        st.session_state.content_data = None
    if "current_step" not in st.session_state:
        st.session_state.current_step = "upload"
    if "image_descriptions" not in st.session_state:
        st.session_state.image_descriptions = []
    if "generated_images" not in st.session_state:
        st.session_state.generated_images = []

    st.title("PDF to PowerPoint Presentation Creator")
    st.write("Transform your PDF documents into professional presentations with AI")

    # Sidebar Configuration
    st.sidebar.title("Presentation Settings")
    
    # Model Selection
    model_deployments = ["gpt4o", "gpt4omini"]
    st.session_state.model_deployment = st.sidebar.selectbox(
        "Select GPT Model",
        model_deployments,
        index=0
    )
    
    # Configuration Options
    st.sidebar.subheader("Configuration")
    num_slides = st.sidebar.number_input("Number of Content Slides", min_value=3, value=5)

    # Visual Options
    st.sidebar.subheader("Visual Settings")
    include_images = st.sidebar.toggle("Generate Images", value=True)
    
    # Special Slides Options
    st.sidebar.subheader("Special Slides")
    include_contents = st.sidebar.toggle("Include Contents Slide", value=True)
    include_conclusion = st.sidebar.toggle("Include Conclusion Slide", value=True)
    include_references = st.sidebar.toggle("Include References Slide", value=True)
    
    # Theme and Layout Settings
    st.sidebar.subheader("Design Settings")
    theme = st.sidebar.selectbox(
        "Select Theme",
        ["Professional", "Modern", "Creative"],
        index=0
    )
    
    layout_style = st.sidebar.selectbox(
        "Select Layout Style",
        ["Content with Image", "Bullet Points", "Minimal"],
        index=0
    )
    
    presentation_style = st.sidebar.selectbox(
        "Presentation Style",
        ["Business", "Academic", "Creative", "Technical"],
        index=0
    )

    # Step 1: PDF Upload
    if st.session_state.current_step == "upload":
        uploaded_file = st.file_uploader("Choose a PDF file", type="pdf")
        if uploaded_file is not None:
            with st.spinner("Processing PDF and generating presentation content..."):
                text_content = extract_text_from_pdf(uploaded_file)
                st.session_state.content_data = get_content_sections(
                    st.session_state.client,
                    text_content,
                    presentation_style,
                    num_slides,
                    include_contents,
                    include_conclusion,
                    include_references
                )
                
                if include_images:
                    st.session_state.current_step = "generate_images"
                else:
                    st.session_state.current_step = "create_presentation"
                st.rerun()

    # Step 2: Edit Content
    elif st.session_state.current_step == "edit_content":
        st.subheader("Edit Presentation Content")
        
        # Title and Subtitle editing
        new_title = st.text_input("Presentation Title", st.session_state.content_data["title"])
        new_subtitle = st.text_input("Presentation Subtitle", st.session_state.content_data["subtitle"])
        
        # Slide content editing
        edited_slides = []
        for idx, slide in enumerate(st.session_state.content_data["slides"]):
            with st.expander(f"Slide {idx + 1}: {slide['title']}", expanded=True):
                new_slide = slide.copy()  # Preserve special slide information
                new_slide["title"] = st.text_input(f"Slide {idx + 1} Title", slide["title"], key=f"title_{idx}")
                new_slide["content"] = st.text_area(f"Slide {idx + 1} Content", slide["content"], key=f"content_{idx}")
                
                key_points = slide.get("key_points", [""])
                new_key_points = []
                for i, point in enumerate(key_points):
                    new_point = st.text_input(f"Key Point {i + 1}", point, key=f"point_{idx}_{i}")
                    new_key_points.append(new_point)
                new_slide["key_points"] = new_key_points
                
                edited_slides.append(new_slide)

        if st.button("Next Step"):
            st.session_state.content_data["title"] = new_title
            st.session_state.content_data["subtitle"] = new_subtitle
            st.session_state.content_data["slides"] = edited_slides
            
            if include_images:
                st.session_state.current_step = "generate_images"
            else:
                st.session_state.current_step = "create_presentation"
            st.rerun()

    # Step 3: Generate Images (if enabled)
    elif st.session_state.current_step == "generate_images":
        st.subheader("Generate Images")
        
        # Generate images for non-special slides
        regular_slides = [slide for slide in st.session_state.content_data["slides"] 
                        if not slide.get("is_special", False)]
        
        total_slides = len(st.session_state.content_data["slides"])
        progress_bar = st.progress(0)
        st.session_state.generated_images = []
        
        for idx, slide in enumerate(st.session_state.content_data["slides"]):
            if slide.get("is_special", False):
                st.session_state.generated_images.append(None)
                continue
                
            with st.spinner(f"Generating image {idx + 1}/{total_slides}..."):
                try:
                    prompt = f"Create a professional presentation image for slide titled '{slide['title']}' with content: {slide['content']}"
                    image = generate_dalle_image(st.session_state.client, prompt)
                    st.session_state.generated_images.append(image)
                except Exception as e:
                    st.warning(f"Failed to generate image for slide {idx + 1}: {str(e)}")
                    st.session_state.generated_images.append(None)
            
            # Calculate progress based on current position relative to total slides
            # Ensure progress never exceeds 1.0
            progress = min(1.0, (idx + 1) / total_slides)
            progress_bar.progress(progress)
        
        st.session_state.current_step = "create_presentation"
        st.rerun()

    # Step 4: Create and Download Presentation
    elif st.session_state.current_step == "create_presentation":
        st.subheader("Create Presentation")
        
        # Combine all data for presentation
        presentation_data = {
            "title": st.session_state.content_data["title"],
            "subtitle": st.session_state.content_data["subtitle"],
            "slides": []
        }
        
        for idx, slide in enumerate(st.session_state.content_data["slides"]):
            slide_data = slide.copy()
            if include_images and idx < len(st.session_state.generated_images):
                slide_data["image"] = st.session_state.generated_images[idx]
            presentation_data["slides"].append(slide_data)

        with st.spinner("Creating presentation..."):
            presentation = create_presentation(
                presentation_data,
                theme,
                layout_style,
                include_images
            )
            
            pptx_io = BytesIO()
            presentation.save(pptx_io)
            pptx_io.seek(0)
            
            st.download_button(
                label="Download Presentation",
                data=pptx_io,
                file_name="generated_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        
        if st.button("Create New Presentation"):
            for key in ["content_data", "current_step", "image_descriptions", "generated_images"]:
                if key in st.session_state:
                    del st.session_state[key]
            st.rerun()

if __name__ == "__main__":
    main()