import streamlit as st
from openai import AzureOpenAI
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import requests
from io import BytesIO
import os
from PIL import Image
import json
from dotenv import load_dotenv

load_dotenv()

def initialize_azure_client():
    """Initialize Azure OpenAI client"""
    client = AzureOpenAI(
        api_key=os.getenv("api_key"),  
        api_version="2024-02-15-preview",
        azure_endpoint=os.getenv("endpoint")
    )
    return client

def extract_text_from_pdf(pdf_file):
    """Extract text from PDF file"""
    text_content = []
    pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
    
    for page_num in range(pdf_document.page_count):
        page = pdf_document[page_num]
        text_content.append(page.get_text())
    
    return " ".join(text_content)

def get_content_sections(client, text, style, max_tokens=2000):
    """Use Azure OpenAI to analyze and structure the content"""
    system_prompt = """You are a professional presentation creator. Your task is to create a presentation outline from the provided text.
    You must respond with valid JSON only, using the following structure:
    {
        "title": "Main presentation title",
        "subtitle": "Subtitle or tagline",
        "slides": [
            {
                "title": "Slide title",
                "content": "Slide content (2-3 sentences)",
                "image_description": "Description for slide image"
            }
        ]
    }
    Do not include any additional text or explanations outside of the JSON structure."""

    user_prompt = f"""Create a {style} presentation structure from the following text. 
    Include engaging titles and clear, concise content for each slide.
    
    Text: {text[:max_tokens]}"""

    try:
        response = client.chat.completions.create(
            model=st.session_state.model_deployment,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7,
            max_tokens=1500,
            response_format={ "type": "json_object" }
        )
        
        content = response.choices[0].message.content
        
        try:
            parsed_content = json.loads(content)
            
            required_fields = ['title', 'subtitle', 'slides']
            if not all(field in parsed_content for field in required_fields):
                raise ValueError("Missing required fields in response")
            
            for slide in parsed_content['slides']:
                if not all(field in slide for field in ['title', 'content', 'image_description']):
                    raise ValueError("Invalid slide structure in response")
            
            return parsed_content
            
        except json.JSONDecodeError as e:
            st.error(f"Error parsing presentation structure. Using default template.")
            return {
                "title": "Document Analysis",
                "subtitle": "Generated Summary",
                "slides": [
                    {
                        "title": "Key Points",
                        "content": text[:200] + "...",
                        "image_description": "Document summary visualization"
                    }
                ]
            }
            
    except Exception as e:
        st.error(f"Error generating presentation structure: {str(e)}")
        return {
            "title": "Document Analysis",
            "subtitle": "Generated Summary",
            "slides": [
                {
                    "title": "Key Points",
                    "content": text[:200] + "...",
                    "image_description": "Document summary visualization"
                }
            ]
        }

def apply_theme_colors(theme):
    """Return extended color scheme based on selected theme"""
    color_schemes = {
        "Professional": {
            "background": RGBColor(255, 255, 255),
            "title": RGBColor(31, 73, 125),
            "text": RGBColor(0, 0, 0),
            "accent": RGBColor(0, 112, 192),
            "footer": RGBColor(89, 89, 89),
            "bullet": RGBColor(31, 73, 125)
        },
        "Modern": {
            "background": RGBColor(240, 240, 240),
            "title": RGBColor(41, 128, 185),
            "text": RGBColor(44, 62, 80),
            "accent": RGBColor(52, 152, 219),
            "footer": RGBColor(149, 165, 166),
            "bullet": RGBColor(41, 128, 185)
        },
        "Creative": {
            "background": RGBColor(255, 250, 240),
            "title": RGBColor(230, 126, 34),
            "text": RGBColor(44, 62, 80),
            "accent": RGBColor(243, 156, 18),
            "footer": RGBColor(149, 165, 166),
            "bullet": RGBColor(230, 126, 34)
        }
    }
    return color_schemes.get(theme, color_schemes["Professional"])

def add_slide_number(slide, slide_number, total_slides, colors):
    """Add slide number to the slide"""
    footer_shape = slide.shapes.add_textbox(
        left=Inches(12),
        top=Inches(7),
        width=Inches(1),
        height=Inches(0.3)
    )
    footer_text_frame = footer_shape.text_frame
    p = footer_text_frame.paragraphs[0]
    p.text = f"{slide_number}/{total_slides}"
    p.font.size = Pt(12)
    p.font.name = "Arial"
    p.font.color.rgb = colors["footer"]
    p.alignment = PP_ALIGN.RIGHT

def apply_content_formatting(paragraph, colors, is_bullet=False):
    """Apply consistent formatting to content paragraphs"""
    paragraph.font.size = Pt(14)
    paragraph.font.name = "Arial"
    paragraph.font.color.rgb = colors["text"]
    
    if is_bullet:
        paragraph.level = 0
        paragraph.bullet = True
        # Format bullet points
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.font.color.rgb = colors["bullet"]
    
    # Add proper line spacing
    paragraph.line_spacing = 1.15
    # Add space after paragraph
    paragraph.space_after = Pt(12)

def create_presentation(content_data, theme, layout_style):
    """Create PowerPoint presentation using python-pptx with enhanced formatting"""
    prs = Presentation()
    colors = apply_theme_colors(theme)
    
    # Set slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    total_slides = len(content_data["slides"]) + 1
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Format title
    title = slide.shapes.title
    title.top = Inches(2)
    title.left = Inches(1)
    title.width = Inches(11.333)
    title.height = Inches(1.5)
    
    # Format subtitle
    subtitle = slide.placeholders[1]
    subtitle.top = Inches(3.75)
    subtitle.left = Inches(1)
    subtitle.width = Inches(11.333)
    
    # Set title and subtitle text with formatting
    title.text = content_data["title"]
    title_para = title.text_frame.paragraphs[0]
    title_para.font.name = "Arial"
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = colors["title"]
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle.text = content_data["subtitle"]
    subtitle_para = subtitle.text_frame.paragraphs[0]
    subtitle_para.font.name = "Arial"
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = colors["text"]
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add slide number
    add_slide_number(slide, 1, total_slides, colors)
    
    # Content slides
    for idx, slide_data in enumerate(content_data["slides"], start=1):
        # Choose layout based on style
        layout = prs.slide_layouts[1]  # Default to title and content
        if layout_style == "Content with Image":
            layout = prs.slide_layouts[8]
        
        slide = prs.slides.add_slide(layout)
        
        # Add title with formatting
        title_shape = slide.shapes.title
        title_shape.top = Inches(0.4)
        title_shape.left = Inches(0.5)
        title_shape.width = Inches(12.33)
        title_shape.height = Inches(0.8)
        
        title_shape.text = slide_data["title"]
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.name = "Arial"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = colors["title"]
        title_para.alignment = PP_ALIGN.LEFT
        
        # Add accent line under title
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            Inches(1.3),
            Inches(12.33),
            Inches(0.03)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = colors["accent"]
        accent_line.line.fill.background()
        
        # Add content
        content_box = slide.shapes.add_textbox(
            left=Inches(0.7),
            top=Inches(1.4),
            width=Inches(11.93),
            height=Inches(5.2)
        )
        
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        # Split content into paragraphs
        content_lines = slide_data["content"].split('. ')
        
        for i, line in enumerate(content_lines):
            if line:
                p = text_frame.add_paragraph()
                p.text = line.strip() + ('.' if not line.endswith('.') else '')
                p.font.name = "Arial"
                p.font.size = Pt(14)
                p.font.color.rgb = colors["text"]
                
                if layout_style == "Bullet Points":
                    p.level = 0
                    p.bullet = True
                
                # Set line spacing
                p.line_spacing = 1.15
                p.space_after = Pt(12)
        
        # Add slide number
        add_slide_number(slide, idx + 1, total_slides, colors)
    
    return prs

def main():
    st.set_page_config(page_title="PDF to PowerPoint Creator", layout="wide")
    
    if "client" not in st.session_state:
        st.session_state.client = initialize_azure_client()
    
    st.sidebar.title("Presentation Settings")
    
    model_deployments = ["gpt4o", "gpt4omini"]
    st.session_state.model_deployment = st.sidebar.selectbox(
        "Select GPT Model",
        model_deployments,
        index=0
    )
    
    theme = st.sidebar.selectbox(
        "Select Theme",
        ["Professional", "Modern", "Creative"],
        index=0
    )
    
    layout_style = st.sidebar.selectbox(
        "Select Layout Style",
        ["Content with Image", "Bullet Points", "Minimal"],
        index=0
    )
    
    presentation_style = st.sidebar.selectbox(
        "Presentation Style",
        ["Business", "Academic", "Creative", "Technical"],
        index=0
    )
    
    st.title("PDF to PowerPoint Presentation Creator")
    st.write("Transform your PDF documents into professional presentations with AI")
    
    uploaded_file = st.file_uploader("Choose a PDF file", type="pdf")
    
    if uploaded_file is not None:
        col1, col2 = st.columns([2, 1])
        
        with col1:
            with st.spinner("Processing PDF..."):
                text_content = extract_text_from_pdf(uploaded_file)
                content_data = get_content_sections(
                    st.session_state.client,
                    text_content,
                    presentation_style
                )
                
                presentation = create_presentation(content_data, theme, layout_style)
                
                pptx_io = BytesIO()
                presentation.save(pptx_io)
                pptx_io.seek(0)
                
                st.download_button(
                    label="Download Presentation",
                    data=pptx_io,
                    file_name="generated_presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
        
        with col2:
            st.subheader("Content Preview")
            st.json(content_data)
            
            st.subheader("Additional Options")
            include_table_contents = st.checkbox("Include Table of Contents", value=True)
            include_sources = st.checkbox("Include Sources Slide", value=True)
            auto_generate_images = st.checkbox("Auto-generate Images", value=True)

if __name__ == "__main__":
    main()